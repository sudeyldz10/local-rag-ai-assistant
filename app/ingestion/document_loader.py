import os
import fitz
from docx import Document
import pytesseract
from PIL import Image


def load_txt(file_path):
    # Plain text files, just read as-is
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()
    return text


def load_pdf(file_path):
    # Extract text page by page using PyMuPDF (fitz)
    text = ""
    pdf = fitz.open(file_path)
    for page in pdf:
        page_text = page.get_text()
        # Drop short/noisy lines (headers, page numbers, stray characters)
        lines = [line.strip() for line in page_text.split("\n") if len(line.strip()) > 20]
        text += "\n".join(lines) + "\n"
    return text


def load_docx(file_path):
    # Word documents: only keep non-empty paragraphs
    doc = Document(file_path)
    parapraghs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(parapraghs)
    return text


def load_png(file_path):
    # Real raster image OCR via pytesseract (fitz has no text layer for actual photos/screenshots)
    text = ""
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang="eng+tur")
    except Exception as e:
        print(f"OCR error for {file_path}: {e}")
        text = ""
    return text.strip()


def load_jpg(file_path):
    # Same handling as PNG
    return load_png(file_path)


def load_jpeg(file_path):
    # Same handling as PNG
    return load_png(file_path)


def load_md(file_path):
    # Markdown is just treated as plain text (no markdown parsing)
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()
    return text


def load_pptx(file_path):
    # Pull text out of every shape on every slide
    from pptx import Presentation
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + ""
    return text


def load_xlsx(file_path):
    # Flatten every cell value across every sheet into one text blob
    import openpyxl
    wb = openpyxl.load_workbook(file_path)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text += str(cell) + " "
    return text


def load_single_file(file_path):
    # Dispatch to the right loader based on file extension
    if file_path.endswith(".txt"):
        return load_txt(file_path)
    elif file_path.endswith(".pdf"):
        return load_pdf(file_path)
    elif file_path.endswith(".docx"):
        return load_docx(file_path)
    elif file_path.endswith(".png"):
        return load_png(file_path)
    elif file_path.endswith(".jpg") or file_path.endswith(".jpeg"):
        return load_jpg(file_path) or load_jpeg(file_path)
    elif file_path.endswith(".md"):
        return load_md(file_path)
    elif file_path.endswith(".pptx"):
        return load_pptx(file_path)
    elif file_path.endswith(".xlsx"):
        return load_xlsx(file_path)
    else:
        print(f"unsupported formats: {file_path}")
        return None


def load_documents(data_dir="data"):
    # Walk the whole data directory (including subfolders) and load every supported file
    supported = {".txt", ".pdf", ".docx", ".md", ".png", ".jpg", ".jpeg", ".pptx", ".xlsx"}
    all_docs = []

    for root, dirs, files in os.walk(data_dir):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            ext = os.path.splitext(file_name)[1].lower()

            if ext not in supported:
                print(f"unsuppoted formats: {file_path}")
                continue

            try:
                text = load_single_file(file_path)
                if text:
                    all_docs.append({"text": text, "source": file_path})
                    print(f"Downloaded: {file_path}")
            except Exception as e:
                # Don't let one bad file kill the whole ingestion run
                print(f"Error: {file_name} - {e}")

    print(f"\nTotal {len(all_docs)} is downloaded")
    return all_docs